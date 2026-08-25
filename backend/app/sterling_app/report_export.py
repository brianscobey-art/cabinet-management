"""Render a Sterling report to .xlsx or .pptx.

Both readers take the same {"report", "meta", "rows"} structure that reports.run()
produces, so a column added to a report definition appears in both exports with
no change here.
"""

from __future__ import annotations

import io

from app.sterling_app.reports import CARTER_GREEN, NEG, POS, WARN

MONEY_FMT = '$#,##0.00;[Red]-$#,##0.00'
PCT_FMT = '0.0%'


# ---------------------------------------------------------------- Excel -----
def _sheet(wb, title, cols, rows, *, heading=None, subtitle=None, notes=None,
           first=False):
    """One formatted sheet: heading, header band, rows, autofilter, notes."""
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    ws = wb.active if first else wb.create_sheet()
    # Excel caps sheet names at 31 chars and forbids : \ / ? * [ ]
    safe = title[:31]
    for ch in (":", chr(92), "/", "?", "*", "[", "]"):
        safe = safe.replace(ch, "-")
    ws.title = safe

    head = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    fill = PatternFill("solid", fgColor=CARTER_GREEN)
    thin = Side(style="thin", color="D9D9D9")

    r = 1
    if heading:
        ws.cell(row=1, column=1, value=heading).font = Font(
            name="Calibri", size=15, bold=True, color=CARTER_GREEN)
        r = 2
        if subtitle:
            ws.cell(row=2, column=1, value=subtitle).font = Font(
                name="Calibri", size=10, color="666666")
            r = 3
        r += 1

    header_row = r
    for i, col in enumerate(cols, start=1):
        cell = ws.cell(row=header_row, column=i, value=col["label"])
        cell.font = head
        cell.fill = fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    status_key = "action" if any(c["key"] == "action" for c in cols) else "status"
    for j, row in enumerate(rows, start=header_row + 1):
        for i, col in enumerate(cols, start=1):
            cell = ws.cell(row=j, column=i, value=row.get(col["key"]))
            cell.border = Border(bottom=thin)
            if col["kind"] == "money":
                cell.number_format = MONEY_FMT
            elif col["kind"] == "pct":
                cell.number_format = PCT_FMT
            if col["kind"] in ("money", "pct", "num"):
                cell.alignment = Alignment(horizontal="right")
        state = str(row.get(status_key, ""))
        if state:
            colour = (NEG if state.startswith(("Below cost", "Adjust"))
                      else WARN if state.startswith(("Below", "Watch")) else POS)
            idx = next(i for i, c in enumerate(cols, start=1) if c["key"] == status_key)
            ws.cell(row=j, column=idx).font = Font(name="Calibri", size=11, bold=True, color=colour)

    widths = {"text": 26, "num": 11, "money": 14, "pct": 11, "pill": 15}
    for i, col in enumerate(cols, start=1):
        ws.column_dimensions[get_column_letter(i)].width = widths.get(col["kind"], 14)
    ws.freeze_panes = ws.cell(row=header_row + 1, column=1)
    if rows:
        ws.auto_filter.ref = (f"A{header_row}:"
                              f"{get_column_letter(len(cols))}{header_row + len(rows)}")
    if notes:
        n = header_row + len(rows) + 2
        ws.cell(row=n, column=1, value="How this is measured").font = Font(
            name="Calibri", size=10, bold=True)
        ws.cell(row=n + 1, column=1, value=notes).font = Font(
            name="Calibri", size=9, color="666666")
    return ws


def to_xlsx(data: dict) -> io.BytesIO:
    """Summary tab, every plan, then one tab per division."""
    from openpyxl import Workbook
    from openpyxl.styles import Font

    rep, meta, rows = data["report"], data["meta"], data["rows"]
    groups = data.get("groups") or []
    cols = rep["columns"]
    gcols = rep.get("group_columns") or []

    wb = Workbook()
    period = "  |  ".join(b for b in (meta.get("period"), meta.get("source"),
                                      f"generated {meta.get('generated', '')}") if b)

    # --- Summary: divisions, worst margin first -----------------------------
    if groups:
        ws = _sheet(wb, "Summary", gcols, groups, heading=rep["title"],
                    subtitle=period, first=True)
        n = len(groups) + 6
        ws.cell(row=n, column=1, value="Divisions under 10% need pricing adjusted.").font = Font(
            name="Calibri", size=10, bold=True, color=NEG)
        for i, (label, value, note) in enumerate(meta.get("headline", [])):
            ws.cell(row=n + 2 + i, column=1, value=label).font = Font(name="Calibri", size=10)
            ws.cell(row=n + 2 + i, column=2, value=value).font = Font(
                name="Calibri", size=10, bold=True, color=CARTER_GREEN)
            ws.cell(row=n + 2 + i, column=3, value=note).font = Font(
                name="Calibri", size=9, color="999999")
        _sheet(wb, "All plans", cols, rows, heading="Every plan",
               subtitle=period, notes=rep.get("notes"))
    else:
        _sheet(wb, "Report", cols, rows, heading=rep["title"], subtitle=period,
               notes=rep.get("notes"), first=True)

    # --- one tab per division ----------------------------------------------
    gkey = rep.get("group_by")
    if gkey and groups:
        for g in groups:
            name = g["division"]
            sub = [r for r in rows if (r.get(gkey) or "—") == name]
            _sheet(wb, name, cols, sub, heading=name,
                   subtitle=(f"{g['plans']} plans · {g['houses']} houses · "
                             f"margin {g['margin'] * 100:.1f}% · "
                             f"{g['below10']} under 10% · {g['action']}"))

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ------------------------------------------------------------ PowerPoint ----
def to_pptx(data: dict, top_n: int = 12) -> io.BytesIO:
    """Title slide, headline figures, the worst N as a chart, then the detail.

    Deliberately not the whole table: a 67-row grid is unreadable projected. The
    deck carries the argument, the workbook carries the data.
    """
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu, Inches, Pt

    rep, meta, rows = data["report"], data["meta"], data["rows"]
    green = RGBColor.from_string(CARTER_GREEN)
    neg = RGBColor.from_string(NEG)
    grey = RGBColor.from_string("666666")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def textbox(slide, left, top, width, height, text, size, *, bold=False,
                color=None, align_right=False):
        from pptx.enum.text import PP_ALIGN

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = "Calibri"
        if color is not None:
            p.font.color.rgb = color
        if align_right:
            p.alignment = PP_ALIGN.RIGHT
        return box

    # --- 1. title -----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2),
            rep["title"], 40, bold=True, color=green)
    textbox(s, Inches(0.9), Inches(3.5), Inches(10.5), Inches(1.0), rep["blurb"], 18, color=grey)
    sub = "  |  ".join(b for b in (meta.get("period"), f"{meta.get('plans', 0)} plans",
                                   f"{meta.get('houses', 0)} houses",
                                   f"generated {meta.get('generated', '')}") if b)
    textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.5), sub, 13, color=grey)

    # --- 2. headline figures ------------------------------------------------
    s = prs.slides.add_slide(blank)
    textbox(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.7),
            "The numbers", 30, bold=True, color=green)
    tiles = meta.get("headline", [])
    if tiles:
        w = Inches(11.5 / max(len(tiles), 1))
        for i, (label, value, note) in enumerate(tiles):
            left = Inches(0.9) + Emu(int(w) * i)
            textbox(s, left, Inches(2.0), w, Inches(0.4), label.upper(), 12, color=grey)
            textbox(s, left, Inches(2.4), w, Inches(1.0), value, 34, bold=True, color=green)
            textbox(s, left, Inches(3.4), w, Inches(0.6), note, 11, color=grey)

    # --- 2b. divisions ------------------------------------------------------
    groups = data.get("groups") or []
    if groups:
        s = prs.slides.add_slide(blank)
        textbox(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.6),
                "By division", 30, bold=True, color=green)
        textbox(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.4),
                "Margin is dollar-weighted. Anything under 10% needs pricing adjusted.",
                13, color=grey)
        gcols = [c for c in (rep.get("group_columns") or [])
                 if c["key"] in ("division", "plans", "houses", "avg_po",
                                 "margin", "below10", "action")]
        tbl = s.shapes.add_table(len(groups) + 1, len(gcols), Inches(0.9), Inches(1.8),
                                 Inches(11.5), Inches(0.45 + 0.42 * len(groups))).table
        for i, c in enumerate(gcols):
            cell = tbl.cell(0, i)
            cell.text = c["label"]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor.from_string("FFFFFF")
            cell.fill.solid()
            cell.fill.fore_color.rgb = green
        for r_i, g in enumerate(groups, start=1):
            for c_i, c in enumerate(gcols):
                v = g.get(c["key"])
                if c["kind"] == "money":
                    txt = ("-$" if float(v) < 0 else "$") + f"{abs(float(v)):,.0f}"
                elif c["kind"] == "pct":
                    txt = f"{float(v) * 100:.1f}%"
                else:
                    txt = str(v)
                cell = tbl.cell(r_i, c_i)
                cell.text = txt
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                if c["key"] in ("margin", "action") and g["margin"] < 0.10:
                    para.font.bold = True
                    para.font.color.rgb = neg

    # --- 3. worst offenders chart ------------------------------------------
    worst = [r for r in rows if float(r.get("exp", 0)) < 0][:top_n]
    if worst:
        s = prs.slides.add_slide(blank)
        textbox(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.6),
                "Where the money is", 30, bold=True, color=green)
        textbox(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.4),
                "12-month exposure by plan — volume times gap, not margin percentage",
                13, color=grey)
        cd = CategoryChartData()
        cd.categories = [r["plan"] for r in reversed(worst)]
        cd.add_series("Exposure", [abs(float(r["exp"])) for r in reversed(worst)])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(1.7),
                                Inches(11.5), Inches(5.2), cd)
        chart = gf.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 60
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = neg
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.value_axis.tick_labels.font.size = Pt(11)
        chart.value_axis.has_major_gridlines = True

    # --- 4. detail table (top N only) --------------------------------------
    if worst:
        s = prs.slides.add_slide(blank)
        textbox(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.6),
                "Plans to review", 30, bold=True, color=green)
        show = [c for c in rep["columns"]
                if c["key"] in ("plan", "n", "po_price", "avg", "sale", "gap", "exp", "margin")]
        tbl = s.shapes.add_table(len(worst) + 1, len(show), Inches(0.9), Inches(1.4),
                                 Inches(11.5), Inches(0.4 + 0.32 * len(worst))).table
        for i, c in enumerate(show):
            cell = tbl.cell(0, i)
            cell.text = c["label"]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = RGBColor.from_string("FFFFFF")
            cell.fill.solid()
            cell.fill.fore_color.rgb = green
        for r_i, row in enumerate(worst, start=1):
            for c_i, c in enumerate(show):
                v = row.get(c["key"])
                if c["kind"] == "money":
                    txt = ("-$" if float(v) < 0 else "$") + f"{abs(float(v)):,.0f}"
                elif c["kind"] == "pct":
                    txt = f"{float(v) * 100:.1f}%"
                else:
                    txt = str(v)
                cell = tbl.cell(r_i, c_i)
                cell.text = txt
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10.5)
                if c["key"] in ("gap", "exp") and float(row.get(c["key"], 0)) < 0:
                    para.font.color.rgb = neg

    # --- 5. method ----------------------------------------------------------
    if rep.get("notes"):
        s = prs.slides.add_slide(blank)
        textbox(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.7),
                "How this is measured", 28, bold=True, color=green)
        textbox(s, Inches(0.9), Inches(1.6), Inches(11.0), Inches(2.5), rep["notes"], 15, color=grey)
        skipped = meta.get("skipped") or {}
        if skipped:
            txt = "Excluded: " + ", ".join(f"{v} {k}" for k, v in skipped.items())
            textbox(s, Inches(0.9), Inches(4.2), Inches(11.0), Inches(1.0), txt, 13, color=grey)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
